#!/usr/bin/env python3
"""Generate ground truth text files for benchmark fixtures.

Strategy by file type:
- PDF (text layer): pdftotext (poppler)
- PDF (scanned): codex CLI with vision
- DOCX: python-docx paragraph text
- PPTX: python-pptx slide text
- XLSX: openpyxl cell values
- HTML: strip tags, extract text
- MD: raw markdown content
- JPG/PNG/images: codex CLI with vision

Usage:
    python scripts/benchmarks/generate-ground-truth.py [--fixtures-dir DIR] [--dry-run]
"""

import argparse
import json
import os
import re
import subprocess
import sys
from pathlib import Path


FIXTURE_DIR = Path("tools/benchmark-harness/fixtures")
GROUND_TRUTH_DIR = FIXTURE_DIR / "ground_truth"


def extract_pdf_text(doc_path: Path) -> tuple[str, str]:
    """Extract text from PDF using pdftotext."""
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", str(doc_path), "-"],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip(), "pdftotext"
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    return "", "pdftotext"


def extract_docx_text(doc_path: Path) -> tuple[str, str]:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document
        doc = Document(str(doc_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs), "python-docx"
    except ImportError:
        print("Warning: python-docx not installed, skipping DOCX extraction", file=sys.stderr)
        return "", "python-docx"


def extract_pptx_text(doc_path: Path) -> tuple[str, str]:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(str(doc_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts), "python-pptx"
    except ImportError:
        print("Warning: python-pptx not installed, skipping PPTX extraction", file=sys.stderr)
        return "", "python-pptx"


def extract_xlsx_text(doc_path: Path) -> tuple[str, str]:
    """Extract text from XLSX using openpyxl."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(doc_path), read_only=True, data_only=True)
        texts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    texts.append(row_text)
        wb.close()
        return "\n".join(texts), "openpyxl"
    except ImportError:
        print("Warning: openpyxl not installed, skipping XLSX extraction", file=sys.stderr)
        return "", "openpyxl"


def extract_html_text(doc_path: Path) -> tuple[str, str]:
    """Extract text from HTML by stripping tags."""
    content = doc_path.read_text(errors="replace")
    # Simple tag stripping
    text = re.sub(r"<script[^>]*>.*?</script>", "", content, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text, "raw_source"


def extract_markdown_text(doc_path: Path) -> tuple[str, str]:
    """Extract raw markdown content."""
    return doc_path.read_text(errors="replace").strip(), "raw_source"


def extract_image_text(doc_path: Path) -> tuple[str, str]:
    """Extract text from image using codex CLI with vision."""
    try:
        result = subprocess.run(
            [
                "codex", "exec",
                "Extract all text from this image exactly as written, preserving layout. Output plain text only.",
                "--image", str(doc_path),
            ],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip(), "codex-vision"
    except (FileNotFoundError, subprocess.TimeoutExpired):
        print(f"Warning: codex CLI not available for image {doc_path}", file=sys.stderr)
    return "", "codex-vision"


EXTRACTORS = {
    "pdf": extract_pdf_text,
    "docx": extract_docx_text,
    "doc": extract_docx_text,
    "pptx": extract_pptx_text,
    "ppt": extract_pptx_text,
    "xlsx": extract_xlsx_text,
    "xls": extract_xlsx_text,
    "html": extract_html_text,
    "md": extract_markdown_text,
    "jpg": extract_image_text,
    "jpeg": extract_image_text,
    "png": extract_image_text,
    "gif": extract_image_text,
    "bmp": extract_image_text,
    "tiff": extract_image_text,
    "webp": extract_image_text,
}


def process_fixture(fixture_path: Path, dry_run: bool = False) -> bool:
    """Process a single fixture JSON and generate ground truth."""
    with open(fixture_path) as f:
        fixture = json.load(f)

    file_type = fixture.get("file_type", "")
    doc_rel_path = fixture.get("document", "")
    fixture_name = fixture_path.stem

    # Resolve document path relative to fixture file
    doc_path = fixture_path.parent / doc_rel_path

    if not doc_path.exists():
        print(f"  Skipping {fixture_name}: document not found at {doc_path}", file=sys.stderr)
        return False

    extractor = EXTRACTORS.get(file_type.lower())
    if extractor is None:
        print(f"  Skipping {fixture_name}: no extractor for file type '{file_type}'", file=sys.stderr)
        return False

    if dry_run:
        print(f"  Would generate ground truth for {fixture_name} ({file_type})")
        return True

    text, source = extractor(doc_path)
    if not text:
        print(f"  Warning: empty extraction for {fixture_name}", file=sys.stderr)
        return False

    # Write ground truth text file
    gt_filename = f"{fixture_name}.txt"
    gt_path = GROUND_TRUTH_DIR / gt_filename
    gt_path.write_text(text)
    print(f"  Generated: {gt_path} ({len(text)} chars, source={source})")

    # Update fixture JSON with ground_truth field
    fixture["ground_truth"] = {
        "text_file": f"ground_truth/{gt_filename}",
        "source": source,
    }
    with open(fixture_path, "w") as f:
        json.dump(fixture, f, indent=2)
        f.write("\n")

    return True


def main():
    parser = argparse.ArgumentParser(description="Generate ground truth for benchmark fixtures")
    parser.add_argument("--fixtures-dir", type=Path, default=FIXTURE_DIR)
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args()

    fixtures_dir = args.fixtures_dir
    gt_dir = fixtures_dir / "ground_truth"

    if not args.dry_run:
        gt_dir.mkdir(parents=True, exist_ok=True)

    fixture_files = sorted(fixtures_dir.glob("*.json"))
    print(f"Found {len(fixture_files)} fixture files in {fixtures_dir}")

    success = 0
    skipped = 0
    for fixture_path in fixture_files:
        print(f"Processing {fixture_path.name}...")
        if process_fixture(fixture_path, dry_run=args.dry_run):
            success += 1
        else:
            skipped += 1

    print(f"\nDone: {success} generated, {skipped} skipped")


if __name__ == "__main__":
    main()
